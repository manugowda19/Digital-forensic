import tempfile
import io
from ctypes import cast, POINTER

from PySide6.QtCore import Qt, QUrl, Slot, QSize
from PySide6.QtGui import QIcon, QPixmap, QImage, QAction, QPageLayout
from PySide6.QtGui import QTransform
from PySide6.QtMultimedia import QMediaPlayer, QAudioOutput
from PySide6.QtMultimediaWidgets import QVideoWidget
from PySide6.QtPrintSupport import QPrinter, QPrintDialog
from PySide6.QtWidgets import (QToolBar, QMessageBox, QScrollArea, QLineEdit, QFileDialog)
from PySide6.QtWidgets import QWidget, QVBoxLayout, QPushButton, QSlider, QLabel, QHBoxLayout, QComboBox, \
    QSpacerItem, QSizePolicy
# from comtypes import CLSCTX_ALL
from fitz import open as fitz_open, Matrix
# from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
import os

if os.name == "nt":  # Windows
    from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume
    from comtypes import CLSCTX_ALL


class UnifiedViewer(QWidget):
    def __init__(self, parent=None):
        super(UnifiedViewer, self).__init__(parent)

        self.layout = QVBoxLayout(self)

        # Initialize the viewers
        self.pdf_viewer = PDFViewer()
        self.picture_viewer = PictureViewer(self)
        self.audio_video_viewer = AudioVideoViewer(self)

        # Add the viewers to the layout
        self.layout.addWidget(self.pdf_viewer)
        self.layout.addWidget(self.picture_viewer)
        self.layout.addWidget(self.audio_video_viewer)

        # Hide all viewers initially
        self.pdf_viewer.hide()
        self.picture_viewer.hide()
        self.audio_video_viewer.hide()

    def load(self, content, file_type="text", file_extension=".txt"):
        # Clear all views first
        self.pdf_viewer.clear()
        self.picture_viewer.clear()
        self.audio_video_viewer.clear()

        # Determine content type and show the appropriate viewer
        if file_type == "text":
            if content.startswith(b"%PDF"):
                self.picture_viewer.hide()
                self.audio_video_viewer.hide()
                self.pdf_viewer.show()
                self.pdf_viewer.display(content)
            else:
                self.pdf_viewer.hide()
                self.audio_video_viewer.hide()
                self.picture_viewer.show()
                self.picture_viewer.display(content)
        elif file_type == "audio" or file_type == "video":
            self.pdf_viewer.hide()
            self.picture_viewer.hide()
            self.audio_video_viewer.show()

            with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
                tmp_file.write(content)
                temp_file_path = tmp_file.name  # Save the temporary file path

                # Make sure to display the correct viewer and pass the file path
            self.audio_video_viewer.display(temp_file_path)

    def display_application_content(self, file_content, full_file_path):
        file_extension = os.path.splitext(full_file_path)[-1].lower()
        file_type = "text"  # default

        audio_extensions = ['.mp3', '.wav', '.aac', '.ogg', '.m4a']
        video_extensions = ['.mp4', '.mkv', '.flv', '.avi', '.mov']
        pdf_extensions = ['.pdf']
        
        # Office document extensions
        excel_extensions = ['.xls', '.xlsx', '.xlsm', '.xlsb']
        word_extensions = ['.doc', '.docx']
        powerpoint_extensions = ['.ppt', '.pptx', '.pptm']

        if file_extension in audio_extensions:
            file_type = "audio"
        elif file_extension in video_extensions:
            file_type = "video"
        elif file_extension in pdf_extensions:
            # PDF files - show in PDF viewer
            self.picture_viewer.hide()
            self.audio_video_viewer.hide()
            self.pdf_viewer.show()
            self.pdf_viewer.display(file_content)
            return
        elif file_extension in excel_extensions:
            # Convert Excel to PDF and display
            self._display_excel_as_pdf(file_content, file_extension)
            return
        elif file_extension in word_extensions:
            # Convert Word to PDF and display
            self._display_word_as_pdf(file_content, file_extension)
            return
        elif file_extension in powerpoint_extensions:
            # Convert PowerPoint to PDF and display
            self._display_powerpoint_as_pdf(file_content, file_extension)
            return
        
        # For other files (images, etc.), use the existing load method
        self.load(file_content, file_type=file_type, file_extension=file_extension)

    def _display_excel_as_pdf(self, file_content, file_extension):
        """Convert Excel file to PDF and display it with recovery options."""
        try:
            # Hide other viewers
            self.picture_viewer.hide()
            self.audio_video_viewer.hide()
            self.pdf_viewer.show()
            
            # Try multiple recovery methods
            pdf_content = self._try_excel_recovery(file_content, file_extension)
            
            if pdf_content:
                self.pdf_viewer.display(pdf_content)
            else:
                from PySide6.QtWidgets import QMessageBox
                QMessageBox.warning(self, "Error", "Failed to load Excel file - file may be corrupted or unsupported format")
        except Exception as e:
            print(f"Error displaying Excel file: {e}")
            from PySide6.QtWidgets import QMessageBox
            QMessageBox.warning(self, "Error", f"Failed to display Excel file: {e}")

    def _try_excel_recovery(self, file_content, file_extension):
        """Try multiple methods to recover and convert Excel file."""
        # Validate file content first
        if not file_content or len(file_content) < 100:
            print("Excel file content is too small or empty")
            return None
        
        # Check if file is actually a ZIP archive (for .xlsx)
        if file_extension in ['.xlsx', '.xlsm', '.xlsb']:
            if not file_content.startswith(b'PK'):  # ZIP files start with PK
                print(f"File does not appear to be a valid {file_extension} file (missing ZIP signature)")
                # Try to extract what we can anyway
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            tmp_file.write(file_content)
            temp_file_path = tmp_file.name
        
        try:
            # Method 1: Try normal loading with openpyxl/xlrd
            pdf_content = self._convert_excel_to_pdf_normal(temp_file_path, file_extension)
            if pdf_content:
                return pdf_content
            
            # Method 2: Try with read_only mode (faster, more tolerant)
            pdf_content = self._convert_excel_to_pdf_readonly(temp_file_path, file_extension)
            if pdf_content:
                return pdf_content
            
            # Method 3: Try pandas as alternative
            pdf_content = self._convert_excel_with_pandas(temp_file_path, file_extension)
            if pdf_content:
                return pdf_content
            
            # Method 4: Try to extract raw data from XML (for .xlsx) - only if it's a ZIP
            if file_extension in ['.xlsx', '.xlsm'] and file_content.startswith(b'PK'):
                pdf_content = self._extract_excel_from_xml(file_content)
                if pdf_content:
                    return pdf_content
            
            # Method 5: Try to extract text from file bytes as last resort
            pdf_content = self._extract_excel_text_raw(file_content, file_extension)
            if pdf_content:
                return pdf_content
            
        finally:
            # Cleanup temp file
            try:
                os.unlink(temp_file_path)
            except:
                pass
        
        return None

    def _convert_excel_to_pdf_normal(self, file_path, file_extension):
        """Normal Excel to PDF conversion."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            # Validate file exists and has content
            if not os.path.exists(file_path) or os.path.getsize(file_path) < 100:
                print("Excel file is too small or doesn't exist")
                return None
            
            # Load workbook - handle both .xlsx and .xls formats
            if file_extension in ['.xlsx', '.xlsm', '.xlsb']:
                import openpyxl
                from openpyxl import load_workbook
                try:
                    wb = load_workbook(file_path, data_only=True)
                    sheet_names = wb.sheetnames
                    get_sheet = lambda name: wb[name]
                    iter_rows = lambda ws: ws.iter_rows(values_only=True)
                except Exception as e:
                    print(f"Failed to load {file_extension} file with openpyxl: {e}")
                    return None
            elif file_extension == '.xls':
                import xlrd
                try:
                    wb = xlrd.open_workbook(file_path)
                    sheet_names = wb.sheet_names()
                    get_sheet = lambda name: wb.sheet_by_name(name)
                    iter_rows = lambda ws: ws.get_rows()
                except Exception as e:
                    print(f"Failed to load .xls file with xlrd: {e}")
                    return None
            else:
                return None
            
            return self._create_pdf_from_excel_data(sheet_names, get_sheet, iter_rows)
        except Exception as e:
            print(f"Normal Excel conversion failed: {e}")
            return None

    def _convert_excel_to_pdf_readonly(self, file_path, file_extension):
        """Try Excel conversion with read_only mode (more tolerant of corruption)."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            if file_extension in ['.xlsx', '.xlsm', '.xlsb']:
                import openpyxl
                from openpyxl import load_workbook
                # Try read_only mode which is more tolerant
                wb = load_workbook(file_path, read_only=True, data_only=True)
                sheet_names = wb.sheetnames
                get_sheet = lambda name: wb[name]
                iter_rows = lambda ws: ws.iter_rows(values_only=True)
            else:
                return None
            
            return self._create_pdf_from_excel_data(sheet_names, get_sheet, iter_rows)
        except Exception as e:
            print(f"Read-only Excel conversion failed: {e}")
            return None

    def _create_pdf_from_excel_data(self, sheet_names, get_sheet, iter_rows):
        """Create PDF from Excel data."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
            pdf_path = pdf_buffer.name
            pdf_buffer.close()
            
            c = canvas.Canvas(pdf_path, pagesize=letter)
            width, height = letter
            y_position = height - 50
            line_height = 20
            margin = 50
            
            # Process each worksheet
            for sheet_name in sheet_names:
                try:
                    ws = get_sheet(sheet_name)
                    
                    # Add sheet name as header
                    c.setFont("Helvetica-Bold", 14)
                    c.drawString(margin, y_position, f"Sheet: {sheet_name}")
                    y_position -= line_height * 2
                    
                    # Process rows
                    row_count = 0
                    for row_idx, row in enumerate(iter_rows(ws), 1):
                        if y_position < margin + 50:  # New page if needed
                            c.showPage()
                            y_position = height - 50
                        
                        x_position = margin
                        for col_idx, cell in enumerate(row, 1):
                            try:
                                # Handle both openpyxl and xlrd cell formats
                                if hasattr(cell, 'value'):
                                    cell_value = cell.value
                                else:
                                    cell_value = cell
                                
                                if cell_value is not None and str(cell_value).strip():
                                    cell_text = str(cell_value)[:50]  # Limit text length
                                    c.setFont("Helvetica", 10)
                                    c.drawString(x_position, y_position, cell_text)
                                    x_position += 100  # Column width
                            except:
                                continue
                        
                        y_position -= line_height
                        row_count += 1
                        if row_count > 1000:  # Limit rows to avoid memory issues
                            break
                    
                    c.showPage()
                    y_position = height - 50
                except Exception as e:
                    print(f"Error processing sheet {sheet_name}: {e}")
                    continue
            
            c.save()
            
            # Read PDF and return
            with open(pdf_path, 'rb') as f:
                pdf_content = f.read()
            
            os.unlink(pdf_path)
            return pdf_content
        except Exception as e:
            print(f"PDF creation from Excel data failed: {e}")
            return None

    def _extract_excel_text_raw(self, file_content, file_extension):
        """Try to extract raw text from Excel file bytes."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            import re
            
            # Try to extract readable text from file bytes
            text_content = []
            
            if file_extension in ['.xlsx', '.xlsm']:
                # For .xlsx, try to find text in XML structure
                content_str = file_content.decode('utf-8', errors='ignore')
                # Look for text between XML tags
                text_pattern = r'<t[^>]*>([^<]+)</t>'
                matches = re.findall(text_pattern, content_str)
                text_content.extend(matches[:500])  # Limit matches
            elif file_extension == '.xls':
                # For .xls, try to extract readable ASCII text
                content_str = file_content.decode('latin-1', errors='ignore')
                # Look for readable text blocks
                text_blocks = re.findall(r'[A-Za-z0-9\s]{15,}', content_str)
                text_content.extend(text_blocks[:200])  # Limit blocks
            
            if text_content:
                # Create PDF from extracted text
                pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                pdf_path = pdf_buffer.name
                pdf_buffer.close()
                
                c = canvas.Canvas(pdf_path, pagesize=letter)
                width, height = letter
                y_position = height - 50
                line_height = 20
                margin = 50
                
                c.setFont("Helvetica-Bold", 14)
                c.drawString(margin, y_position, "Recovered Excel Text (raw extraction)")
                y_position -= line_height * 2
                
                for text in text_content[:300]:  # Limit text blocks
                    if y_position < margin + 50:
                        c.showPage()
                        y_position = height - 50
                    
                    # Clean and display text
                    clean_text = re.sub(r'[^\x20-\x7E\n]', '', str(text))[:100]
                    if clean_text.strip():
                        c.setFont("Helvetica", 10)
                        c.drawString(margin, y_position, clean_text)
                        y_position -= line_height
                
                c.save()
                
                with open(pdf_path, 'rb') as f:
                    pdf_content = f.read()
                
                os.unlink(pdf_path)
                print("Excel text extracted from raw file bytes")
                return pdf_content
        except Exception as e:
            print(f"Raw text extraction from Excel failed: {e}")
        
        return None

    def _extract_excel_from_xml(self, file_content):
        """Try to extract Excel data from XML structure (for corrupted .xlsx files)."""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            # Validate it's a ZIP file first
            if not file_content.startswith(b'PK'):
                print("File is not a ZIP archive, cannot extract XML")
                return None
            
            # .xlsx files are ZIP archives containing XML
            excel_zip = zipfile.ZipFile(io.BytesIO(file_content), 'r')
            
            # Try to find shared strings and sheet data
            shared_strings = []
            try:
                if 'xl/sharedStrings.xml' in excel_zip.namelist():
                    strings_xml = excel_zip.read('xl/sharedStrings.xml')
                    root = ET.fromstring(strings_xml)
                    for si in root.findall('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}si'):
                        t = si.find('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}t')
                        if t is not None and t.text:
                            shared_strings.append(t.text)
            except:
                pass
            
            # Try to extract sheet data
            sheet_data = []
            for name in excel_zip.namelist():
                if name.startswith('xl/worksheets/sheet') and name.endswith('.xml'):
                    try:
                        sheet_xml = excel_zip.read(name)
                        root = ET.fromstring(sheet_xml)
                        # Extract row data
                        for row in root.findall('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}row'):
                            row_data = []
                            for cell in row.findall('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}c'):
                                v = cell.find('.//{http://schemas.openxmlformats.org/spreadsheetml/2006/main}v')
                                if v is not None and v.text:
                                    try:
                                        idx = int(v.text)
                                        if idx < len(shared_strings):
                                            row_data.append(shared_strings[idx])
                                        else:
                                            row_data.append(v.text)
                                    except:
                                        row_data.append(v.text)
                            if row_data:
                                sheet_data.append(row_data)
                    except:
                        continue
            
            excel_zip.close()
            
            if sheet_data:
                # Create PDF from extracted data
                pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                pdf_path = pdf_buffer.name
                pdf_buffer.close()
                
                c = canvas.Canvas(pdf_path, pagesize=letter)
                width, height = letter
                y_position = height - 50
                line_height = 20
                margin = 50
                
                c.setFont("Helvetica-Bold", 14)
                c.drawString(margin, y_position, "Recovered Excel Data (from XML)")
                y_position -= line_height * 2
                
                for row in sheet_data[:500]:  # Limit rows
                    if y_position < margin + 50:
                        c.showPage()
                        y_position = height - 50
                    
                    x_position = margin
                    for cell_value in row[:10]:  # Limit columns
                        cell_text = str(cell_value)[:40]
                        c.setFont("Helvetica", 9)
                        c.drawString(x_position, y_position, cell_text)
                        x_position += 80
                    
                    y_position -= line_height
                
                c.save()
                
                with open(pdf_path, 'rb') as f:
                    pdf_content = f.read()
                
                os.unlink(pdf_path)
                print("Excel data extracted from XML structure")
                return pdf_content
        except Exception as e:
            print(f"XML extraction from Excel failed: {e}")
        
        return None

    def _convert_excel_with_pandas(self, file_path, file_extension):
        """Try to convert Excel using pandas (alternative method)."""
        try:
            import pandas as pd
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            # Try to read with pandas
            if file_extension in ['.xlsx', '.xlsm', '.xlsb']:
                excel_file = pd.ExcelFile(file_path, engine='openpyxl')
            elif file_extension == '.xls':
                excel_file = pd.ExcelFile(file_path, engine='xlrd')
            else:
                return None
            
            pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
            pdf_path = pdf_buffer.name
            pdf_buffer.close()
            
            c = canvas.Canvas(pdf_path, pagesize=letter)
            width, height = letter
            y_position = height - 50
            line_height = 20
            margin = 50
            
            for sheet_name in excel_file.sheet_names:
                try:
                    df = pd.read_excel(excel_file, sheet_name=sheet_name, nrows=500)  # Limit rows
                    
                    c.setFont("Helvetica-Bold", 14)
                    c.drawString(margin, y_position, f"Sheet: {sheet_name}")
                    y_position -= line_height * 2
                    
                    # Write headers
                    if not df.empty:
                        headers = list(df.columns)[:10]  # Limit columns
                        x_position = margin
                        for header in headers:
                            c.setFont("Helvetica-Bold", 10)
                            c.drawString(x_position, y_position, str(header)[:30])
                            x_position += 100
                        y_position -= line_height
                    
                    # Write data
                    for idx, row in df.iterrows():
                        if y_position < margin + 50:
                            c.showPage()
                            y_position = height - 50
                        
                        x_position = margin
                        for col in headers:
                            try:
                                value = str(row[col])[:30] if pd.notna(row[col]) else ""
                                c.setFont("Helvetica", 9)
                                c.drawString(x_position, y_position, value)
                                x_position += 100
                            except:
                                x_position += 100
                        
                        y_position -= line_height
                    
                    c.showPage()
                    y_position = height - 50
                except Exception as e:
                    print(f"Error processing sheet {sheet_name} with pandas: {e}")
                    continue
            
            c.save()
            
            with open(pdf_path, 'rb') as f:
                pdf_content = f.read()
            
            os.unlink(pdf_path)
            print("Excel converted using pandas")
            return pdf_content
        except Exception as e:
            print(f"Pandas Excel conversion failed: {e}")
            return None

    def _display_word_as_pdf(self, file_content, file_extension):
        """Convert Word document to PDF and display it with recovery options."""
        try:
            # Hide other viewers
            self.picture_viewer.hide()
            self.audio_video_viewer.hide()
            self.pdf_viewer.show()
            
            # Try multiple recovery methods
            pdf_content = self._try_word_recovery(file_content, file_extension)
            
            if pdf_content:
                self.pdf_viewer.display(pdf_content)
            else:
                from PySide6.QtWidgets import QMessageBox
                QMessageBox.warning(self, "Error", "Failed to load Word file - file may be corrupted or unsupported format")
        except Exception as e:
            print(f"Error displaying Word file: {e}")
            from PySide6.QtWidgets import QMessageBox
            QMessageBox.warning(self, "Error", f"Failed to display Word file: {e}")

    def _try_word_recovery(self, file_content, file_extension):
        """Try multiple methods to recover and convert Word file."""
        # Validate file content first
        if not file_content or len(file_content) < 100:
            print("Word file content is too small or empty")
            return None
        
        # Check if file is actually a ZIP archive (for .docx)
        if file_extension == '.docx':
            if not file_content.startswith(b'PK'):  # ZIP files start with PK
                print("File does not appear to be a valid .docx file (missing ZIP signature)")
                # Try raw text extraction instead
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            tmp_file.write(file_content)
            temp_file_path = tmp_file.name
        
        try:
            # Method 1: Try normal loading with python-docx
            pdf_content = self._convert_word_to_pdf_normal(temp_file_path)
            if pdf_content:
                return pdf_content
            
            # Method 2: Try to extract from XML structure (for .docx) - only if it's a ZIP
            if file_extension == '.docx' and file_content.startswith(b'PK'):
                pdf_content = self._extract_word_from_xml(file_content)
                if pdf_content:
                    return pdf_content
            
            # Method 3: Try to extract raw text from file
            pdf_content = self._extract_word_text_raw(file_content, file_extension)
            if pdf_content:
                return pdf_content
            
        finally:
            # Cleanup temp file
            try:
                os.unlink(temp_file_path)
            except:
                pass
        
        return None

    def _convert_word_to_pdf_normal(self, file_path):
        """Normal Word to PDF conversion."""
        try:
            from docx import Document as DocxDocument
            from docx.opc.exceptions import PackageNotFoundError
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            # Validate file exists and has content
            if not os.path.exists(file_path) or os.path.getsize(file_path) < 100:
                print("Word file is too small or doesn't exist")
                return None
            
            # Load document
            try:
                doc = DocxDocument(file_path)
            except PackageNotFoundError as e:
                print(f"Word file package not found (corrupted or invalid format): {e}")
                return None
            except Exception as e:
                print(f"Failed to load Word document: {e}")
                return None
            
            return self._create_pdf_from_word_document(doc)
        except Exception as e:
            print(f"Normal Word conversion failed: {e}")
            return None

    def _create_pdf_from_word_document(self, doc):
        """Create PDF from Word document."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
            pdf_path = pdf_buffer.name
            pdf_buffer.close()
            
            c = canvas.Canvas(pdf_path, pagesize=letter)
            width, height = letter
            y_position = height - 50
            line_height = 20
            margin = 50
            
            # Process paragraphs
            for paragraph in doc.paragraphs:
                try:
                    if y_position < margin + 50:  # New page if needed
                        c.showPage()
                        y_position = height - 50
                    
                    text = paragraph.text
                    if text.strip():
                        # Determine font based on style
                        try:
                            if paragraph.style.name.startswith('Heading'):
                                font_size = 16
                                font_name = "Helvetica-Bold"
                            else:
                                font_size = 12
                                font_name = "Helvetica"
                        except:
                            font_size = 12
                            font_name = "Helvetica"
                        
                        c.setFont(font_name, font_size)
                        
                        # Wrap text if needed
                        words = text.split()
                        line = ""
                        for word in words:
                            test_line = line + word + " " if line else word + " "
                            if c.stringWidth(test_line, font_name, font_size) < width - 2 * margin:
                                line = test_line
                            else:
                                if line:
                                    c.drawString(margin, y_position, line.strip())
                                    y_position -= line_height
                                line = word + " "
                        
                        if line:
                            c.drawString(margin, y_position, line.strip())
                            y_position -= line_height
                except Exception as e:
                    print(f"Error processing paragraph: {e}")
                    continue
            
            # Process tables if any
            try:
                for table in doc.tables:
                    if y_position < margin + 50:
                        c.showPage()
                        y_position = height - 50
                    
                    y_position -= line_height
                    for row in table.rows:
                        if y_position < margin + 50:
                            c.showPage()
                            y_position = height - 50
                        
                        x_pos = margin
                        for cell in row.cells:
                            try:
                                cell_text = cell.text[:30]
                                c.setFont("Helvetica", 9)
                                c.drawString(x_pos, y_position, cell_text)
                                x_pos += 120
                            except:
                                x_pos += 120
                        
                        y_position -= line_height
            except:
                pass
            
            c.save()
            
            # Read PDF and return
            with open(pdf_path, 'rb') as f:
                pdf_content = f.read()
            
            os.unlink(pdf_path)
            return pdf_content
        except Exception as e:
            print(f"PDF creation from Word document failed: {e}")
            return None

    def _extract_word_from_xml(self, file_content):
        """Try to extract Word content from XML structure (for corrupted .docx files)."""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            # Validate it's a ZIP file first
            if not file_content.startswith(b'PK'):
                print("File is not a ZIP archive, cannot extract XML")
                return None
            
            # .docx files are ZIP archives containing XML
            docx_zip = zipfile.ZipFile(io.BytesIO(file_content), 'r')
            
            # Try to extract text from document.xml
            text_content = []
            try:
                if 'word/document.xml' in docx_zip.namelist():
                    doc_xml = docx_zip.read('word/document.xml')
                    root = ET.fromstring(doc_xml)
                    
                    # Extract text from paragraphs
                    namespace = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
                    for para in root.findall('.//{}p'.format(namespace)):
                        para_text = []
                        for text_elem in para.findall('.//{}t'.format(namespace)):
                            if text_elem.text:
                                para_text.append(text_elem.text)
                        if para_text:
                            text_content.append(' '.join(para_text))
            except Exception as e:
                print(f"Error extracting from document.xml: {e}")
            
            docx_zip.close()
            
            if text_content:
                # Create PDF from extracted text
                pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                pdf_path = pdf_buffer.name
                pdf_buffer.close()
                
                c = canvas.Canvas(pdf_path, pagesize=letter)
                width, height = letter
                y_position = height - 50
                line_height = 20
                margin = 50
                
                c.setFont("Helvetica-Bold", 14)
                c.drawString(margin, y_position, "Recovered Word Document (from XML)")
                y_position -= line_height * 2
                
                for text in text_content:
                    if y_position < margin + 50:
                        c.showPage()
                        y_position = height - 50
                    
                    # Wrap text
                    words = text.split()
                    line = ""
                    for word in words:
                        test_line = line + word + " " if line else word + " "
                        if c.stringWidth(test_line, "Helvetica", 12) < width - 2 * margin:
                            line = test_line
                        else:
                            if line:
                                c.setFont("Helvetica", 12)
                                c.drawString(margin, y_position, line.strip())
                                y_position -= line_height
                            line = word + " "
                    
                    if line:
                        c.setFont("Helvetica", 12)
                        c.drawString(margin, y_position, line.strip())
                        y_position -= line_height
                
                c.save()
                
                with open(pdf_path, 'rb') as f:
                    pdf_content = f.read()
                
                os.unlink(pdf_path)
                print("Word content extracted from XML structure")
                return pdf_content
        except Exception as e:
            print(f"XML extraction from Word failed: {e}")
        
        return None

    def _extract_word_text_raw(self, file_content, file_extension):
        """Try to extract raw text from Word file bytes."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            import re
            
            # Try to extract readable text from file bytes
            text_content = []
            
            if file_extension == '.docx':
                # For .docx, try to find text in XML structure
                content_str = file_content.decode('utf-8', errors='ignore')
                # Look for text between XML tags
                text_pattern = r'<w:t[^>]*>([^<]+)</w:t>'
                matches = re.findall(text_pattern, content_str)
                text_content.extend(matches)
            elif file_extension == '.doc':
                # For .doc, try to extract readable ASCII text
                content_str = file_content.decode('latin-1', errors='ignore')
                # Look for readable text blocks
                text_blocks = re.findall(r'[A-Za-z0-9\s]{20,}', content_str)
                text_content.extend(text_blocks[:100])  # Limit blocks
            
            if text_content:
                # Create PDF from extracted text
                pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                pdf_path = pdf_buffer.name
                pdf_buffer.close()
                
                c = canvas.Canvas(pdf_path, pagesize=letter)
                width, height = letter
                y_position = height - 50
                line_height = 20
                margin = 50
                
                c.setFont("Helvetica-Bold", 14)
                c.drawString(margin, y_position, "Recovered Word Text (raw extraction)")
                y_position -= line_height * 2
                
                for text in text_content[:200]:  # Limit text blocks
                    if y_position < margin + 50:
                        c.showPage()
                        y_position = height - 50
                    
                    # Clean and display text
                    clean_text = re.sub(r'[^\x20-\x7E\n]', '', str(text))[:200]
                    if clean_text.strip():
                        c.setFont("Helvetica", 11)
                        c.drawString(margin, y_position, clean_text)
                        y_position -= line_height
                
                c.save()
                
                with open(pdf_path, 'rb') as f:
                    pdf_content = f.read()
                
                os.unlink(pdf_path)
                print("Word text extracted from raw file bytes")
                return pdf_content
        except Exception as e:
            print(f"Raw text extraction from Word failed: {e}")
        
        return None

    def _display_powerpoint_as_pdf(self, file_content, file_extension):
        """Convert PowerPoint presentation to PDF and display it with recovery options."""
        try:
            # Hide other viewers
            self.picture_viewer.hide()
            self.audio_video_viewer.hide()
            self.pdf_viewer.show()
            
            # Try multiple recovery methods
            pdf_content = self._try_powerpoint_recovery(file_content, file_extension)
            
            if pdf_content:
                self.pdf_viewer.display(pdf_content)
            else:
                from PySide6.QtWidgets import QMessageBox
                QMessageBox.warning(self, "Error", "Failed to load PowerPoint file - file may be corrupted or unsupported format")
        except Exception as e:
            print(f"Error displaying PowerPoint file: {e}")
            from PySide6.QtWidgets import QMessageBox
            QMessageBox.warning(self, "Error", f"Failed to display PowerPoint file: {e}")

    def _try_powerpoint_recovery(self, file_content, file_extension):
        """Try multiple methods to recover and convert PowerPoint file."""
        # Validate file content first
        if not file_content or len(file_content) < 100:
            print("PowerPoint file content is too small or empty")
            return None
        
        # Check if file is actually a ZIP archive (for .pptx)
        if file_extension in ['.pptx', '.pptm']:
            if not file_content.startswith(b'PK'):  # ZIP files start with PK
                print(f"File does not appear to be a valid {file_extension} file (missing ZIP signature)")
                # Try raw text extraction instead
        
        # Create temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=file_extension) as tmp_file:
            tmp_file.write(file_content)
            temp_file_path = tmp_file.name
        
        try:
            # Method 1: Try normal loading with python-pptx
            pdf_content = self._convert_powerpoint_to_pdf_normal(temp_file_path)
            if pdf_content:
                return pdf_content
            
            # Method 2: Try to extract from XML structure (for .pptx) - only if it's a ZIP
            if file_extension in ['.pptx', '.pptm'] and file_content.startswith(b'PK'):
                pdf_content = self._extract_powerpoint_from_xml(file_content)
                if pdf_content:
                    return pdf_content
            
            # Method 3: Try to extract raw text from file
            pdf_content = self._extract_powerpoint_text_raw(file_content, file_extension)
            if pdf_content:
                return pdf_content
            
        finally:
            # Cleanup temp file
            try:
                os.unlink(temp_file_path)
            except:
                pass
        
        return None

    def _convert_powerpoint_to_pdf_normal(self, file_path):
        """Normal PowerPoint to PDF conversion."""
        try:
            from pptx import Presentation
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            # Load presentation
            prs = Presentation(file_path)
            
            return self._create_pdf_from_powerpoint_presentation(prs)
        except Exception as e:
            print(f"Normal PowerPoint conversion failed: {e}")
            return None

    def _create_pdf_from_powerpoint_presentation(self, prs):
        """Create PDF from PowerPoint presentation."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
            pdf_path = pdf_buffer.name
            pdf_buffer.close()
            
            c = canvas.Canvas(pdf_path, pagesize=letter)
            width, height = letter
            margin = 50
            
            # Process slides
            for slide_num, slide in enumerate(prs.slides, 1):
                try:
                    if slide_num > 1:
                        c.showPage()
                    
                    y_position = height - 50
                    line_height = 20
                    
                    # Add slide number
                    c.setFont("Helvetica-Bold", 12)
                    c.drawString(margin, y_position, f"Slide {slide_num}")
                    y_position -= line_height * 2
                    
                    # Process shapes (text boxes, etc.)
                    for shape in slide.shapes:
                        try:
                            if hasattr(shape, "text") and shape.text:
                                text = shape.text.strip()
                                if text:
                                    # Determine font size based on shape
                                    try:
                                        if hasattr(shape, "font") and shape.font:
                                            font_size = 14
                                        else:
                                            font_size = 12
                                    except:
                                        font_size = 12
                                    
                                    c.setFont("Helvetica", font_size)
                                    
                                    # Wrap and draw text
                                    words = text.split()
                                    line = ""
                                    for word in words:
                                        test_line = line + word + " " if line else word + " "
                                        if c.stringWidth(test_line, "Helvetica", font_size) < width - 2 * margin:
                                            line = test_line
                                        else:
                                            if line:
                                                c.drawString(margin, y_position, line.strip())
                                                y_position -= line_height
                                            line = word + " "
                                    
                                    if line:
                                        c.drawString(margin, y_position, line.strip())
                                        y_position -= line_height
                        except Exception as e:
                            print(f"Error processing shape: {e}")
                            continue
                except Exception as e:
                    print(f"Error processing slide {slide_num}: {e}")
                    continue
            
            c.save()
            
            # Read PDF and return
            with open(pdf_path, 'rb') as f:
                pdf_content = f.read()
            
            os.unlink(pdf_path)
            return pdf_content
        except Exception as e:
            print(f"PDF creation from PowerPoint failed: {e}")
            return None

    def _extract_powerpoint_from_xml(self, file_content):
        """Try to extract PowerPoint content from XML structure (for corrupted .pptx files)."""
        try:
            import zipfile
            import xml.etree.ElementTree as ET
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            
            # Validate it's a ZIP file first
            if not file_content.startswith(b'PK'):
                print("File is not a ZIP archive, cannot extract XML")
                return None
            
            # .pptx files are ZIP archives containing XML
            pptx_zip = zipfile.ZipFile(io.BytesIO(file_content), 'r')
            
            # Try to extract text from slide XML files
            slide_texts = []
            for name in pptx_zip.namelist():
                if name.startswith('ppt/slides/slide') and name.endswith('.xml'):
                    try:
                        slide_xml = pptx_zip.read(name)
                        root = ET.fromstring(slide_xml)
                        
                        # Extract text from slide
                        namespace = '{http://schemas.openxmlformats.org/presentationml/2006/main}'
                        text_namespace = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
                        
                        slide_text = []
                        for text_elem in root.findall('.//{}t'.format(text_namespace)):
                            if text_elem.text:
                                slide_text.append(text_elem.text)
                        
                        if slide_text:
                            slide_texts.append(' '.join(slide_text))
                    except:
                        continue
            
            pptx_zip.close()
            
            if slide_texts:
                # Create PDF from extracted text
                pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                pdf_path = pdf_buffer.name
                pdf_buffer.close()
                
                c = canvas.Canvas(pdf_path, pagesize=letter)
                width, height = letter
                y_position = height - 50
                line_height = 20
                margin = 50
                
                for slide_num, text in enumerate(slide_texts, 1):
                    if slide_num > 1:
                        c.showPage()
                        y_position = height - 50
                    
                    c.setFont("Helvetica-Bold", 12)
                    c.drawString(margin, y_position, f"Slide {slide_num} (Recovered from XML)")
                    y_position -= line_height * 2
                    
                    # Wrap and display text
                    words = text.split()
                    line = ""
                    for word in words:
                        test_line = line + word + " " if line else word + " "
                        if c.stringWidth(test_line, "Helvetica", 12) < width - 2 * margin:
                            line = test_line
                        else:
                            if line:
                                c.setFont("Helvetica", 12)
                                c.drawString(margin, y_position, line.strip())
                                y_position -= line_height
                            line = word + " "
                    
                    if line:
                        c.setFont("Helvetica", 12)
                        c.drawString(margin, y_position, line.strip())
                        y_position -= line_height
                
                c.save()
                
                with open(pdf_path, 'rb') as f:
                    pdf_content = f.read()
                
                os.unlink(pdf_path)
                print("PowerPoint content extracted from XML structure")
                return pdf_content
        except Exception as e:
            print(f"XML extraction from PowerPoint failed: {e}")
        
        return None

    def _extract_powerpoint_text_raw(self, file_content, file_extension):
        """Try to extract raw text from PowerPoint file bytes."""
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas
            import re
            
            # Try to extract readable text from file bytes
            text_content = []
            
            if file_extension in ['.pptx', '.pptm']:
                # For .pptx, try to find text in XML structure
                content_str = file_content.decode('utf-8', errors='ignore')
                # Look for text between XML tags
                text_pattern = r'<a:t[^>]*>([^<]+)</a:t>'
                matches = re.findall(text_pattern, content_str)
                text_content.extend(matches)
            elif file_extension == '.ppt':
                # For .ppt, try to extract readable ASCII text
                content_str = file_content.decode('latin-1', errors='ignore')
                # Look for readable text blocks
                text_blocks = re.findall(r'[A-Za-z0-9\s]{20,}', content_str)
                text_content.extend(text_blocks[:50])  # Limit blocks
            
            if text_content:
                # Create PDF from extracted text
                pdf_buffer = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
                pdf_path = pdf_buffer.name
                pdf_buffer.close()
                
                c = canvas.Canvas(pdf_path, pagesize=letter)
                width, height = letter
                y_position = height - 50
                line_height = 20
                margin = 50
                
                c.setFont("Helvetica-Bold", 14)
                c.drawString(margin, y_position, "Recovered PowerPoint Text (raw extraction)")
                y_position -= line_height * 2
                
                for idx, text in enumerate(text_content[:100], 1):  # Limit text blocks
                    if y_position < margin + 50:
                        c.showPage()
                        y_position = height - 50
                    
                    # Clean and display text
                    clean_text = re.sub(r'[^\x20-\x7E\n]', '', str(text))[:150]
                    if clean_text.strip():
                        c.setFont("Helvetica", 11)
                        c.drawString(margin, y_position, f"{idx}. {clean_text}")
                        y_position -= line_height
                
                c.save()
                
                with open(pdf_path, 'rb') as f:
                    pdf_content = f.read()
                
                os.unlink(pdf_path)
                print("PowerPoint text extracted from raw file bytes")
                return pdf_content
        except Exception as e:
            print(f"Raw text extraction from PowerPoint failed: {e}")
        
        return None

    def clear(self):
        self.pdf_viewer.clear()
        self.picture_viewer.clear()
        self.audio_video_viewer.clear()


class PictureViewer(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.original_pixmap = None  # Store the original QPixmap
        self.original_image_bytes = None  # Store the original image bytes
        self.initialize_ui()

    def initialize_ui(self):
        self.layout = QVBoxLayout()
        self.layout.setContentsMargins(0, 0, 0, 0)
        self.layout.setSpacing(0)
        self.layout.setAlignment(Qt.AlignCenter)

        # Create a container for the toolbar and the application viewer
        container_widget = QWidget(self)

        container_layout = QVBoxLayout()
        container_layout.setContentsMargins(0, 0, 0, 0)  # Remove any margins
        container_layout.setSpacing(0)  # Remove spacing between toolbar and viewer

        # Create and set up the toolbar
        self.setup_toolbar()

        # Add the toolbar to the container layout
        container_layout.addWidget(self.toolbar)

        self.image_label = QLabel(self)
        self.image_label.setContentsMargins(0, 0, 0, 0)
        self.image_label.setAlignment(Qt.AlignCenter)

        self.scroll_area = QScrollArea(self)
        self.scroll_area.setContentsMargins(0, 0, 0, 0)
        self.scroll_area.setWidget(self.image_label)
        self.scroll_area.setWidgetResizable(True)

        container_layout.addWidget(self.scroll_area)
        container_widget.setLayout(container_layout)
        self.layout.addWidget(container_widget)
        self.setLayout(self.layout)

    def setup_toolbar(self):
        self.toolbar = QToolBar(self)
        self.toolbar.setContentsMargins(0, 0, 0, 0)
        # self.toolbar.setStyleSheet("QToolBar { background-color: lightgray; border: 0px solid gray; }")

        zoom_in_icon = QIcon("Icons/icons8-zoom-in-50.png")
        zoom_out_icon = QIcon("Icons/icons8-zoom-out-50.png")
        rotate_left_icon = QIcon("Icons/icons8-rotate-left-50.png")
        rotate_right_icon = QIcon("Icons/icons8-rotate-right-50.png")
        reset_icon = QIcon("Icons/icons8-no-rotation-50.png")
        export_icon = QIcon("Icons/icons8-save-as-50.png")

        zoom_in_action = QAction(zoom_in_icon, 'Zoom In', self)
        zoom_out_action = QAction(zoom_out_icon, 'Zoom Out', self)
        rotate_left_action = QAction(rotate_left_icon, 'Rotate Left', self)
        rotate_right_action = QAction(rotate_right_icon, 'Rotate Right', self)
        reset_action = QAction(reset_icon, 'Reset', self)
        self.export_action = QAction(export_icon, 'Save Image', self)

        zoom_in_action.triggered.connect(self.zoom_in)
        zoom_out_action.triggered.connect(self.zoom_out)
        rotate_left_action.triggered.connect(self.rotate_left)
        rotate_right_action.triggered.connect(self.rotate_right)
        reset_action.triggered.connect(self.reset)
        self.export_action.triggered.connect(self.export_original_image)

        # Add actions to the toolbar
        self.toolbar.addAction(zoom_in_action)
        self.toolbar.addAction(zoom_out_action)
        self.toolbar.addAction(rotate_left_action)
        self.toolbar.addAction(rotate_right_action)
        self.toolbar.addAction(reset_action)
        self.toolbar.addAction(self.export_action)

    def display(self, content):
        self.original_image_bytes = content  # Save the original image bytes
        # Convert byte data to QPixmap
        qt_image = QImage.fromData(content)
        pixmap = QPixmap.fromImage(qt_image)
        self.original_pixmap = pixmap.copy()  # Save the original pixmap
        self.image_label.setPixmap(pixmap)

    def clear(self):
        self.image_label.clear()

    def zoom_in(self):
        self.image_label.setPixmap(self.image_label.pixmap().scaled(
            self.image_label.width() * 1.2, self.image_label.height() * 1.2, Qt.KeepAspectRatio,
            Qt.SmoothTransformation))

    def zoom_out(self):
        self.image_label.setPixmap(self.image_label.pixmap().scaled(
            self.image_label.width() * 0.8, self.image_label.height() * 0.8, Qt.KeepAspectRatio,
            Qt.SmoothTransformation))

    def rotate_left(self):
        transform = QTransform().rotate(-90)
        pixmap = self.image_label.pixmap().transformed(transform)
        self.image_label.setPixmap(pixmap)

    def rotate_right(self):
        transform = QTransform().rotate(90)
        pixmap = self.image_label.pixmap().transformed(transform)
        self.image_label.setPixmap(pixmap)

    def reset(self):
        if self.original_pixmap:
            self.image_label.setPixmap(self.original_pixmap)

    def export_original_image(self):
        # Ensure that an image is currently loaded
        if not self.original_image_bytes:
            QMessageBox.warning(self, "Export Error", "No image is currently loaded.")
            return

        # Ask the user where to save the exported image
        file_name, _ = QFileDialog.getSaveFileName(self, "Export Image", "",
                                                   "PNG (*.png);;JPEG (*.jpg *.jpeg);;All Files (*)")

        # If a location is chosen, save the image
        if file_name:
            with open(file_name, 'wb') as f:
                f.write(self.original_image_bytes)
            QMessageBox.information(self, "Export Success", "Image exported successfully!")


class PDFViewer(QWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.pdf = None
        self.current_page = 0
        self.zoom_factor = 1.0  # Initialize the zoom factor here
        self.rotation_angle = 0
        self.is_panning = False
        self.pan_start_x = 0
        self.pan_start_y = 0
        self.pan_mode = False

        self.initialize_ui()
        if self.pdf:
            self.show_page(self.current_page)

    def initialize_ui(self):
        # Set up the main layout
        self.layout = QVBoxLayout()
        self.layout.setContentsMargins(0, 0, 0, 0)
        self.layout.setSpacing(0)
        self.layout.setAlignment(Qt.AlignCenter)

        # Create a container for the toolbar and the application viewer
        container_widget = QWidget(self)
        container_layout = QVBoxLayout()
        container_layout.setContentsMargins(0, 0, 0, 0)
        container_layout.setSpacing(0)

        # Create and set up the toolbar
        self.setup_toolbar()

        # Add the toolbar to the container layout
        container_layout.addWidget(self.toolbar)

        # Set up the PDF display area
        self.setup_pdf_display_area()
        container_layout.addWidget(self.scroll_area)

        container_widget.setLayout(container_layout)
        self.layout.addWidget(container_widget)

        self.setLayout(self.layout)
        self.update_navigation_states()

    def setup_toolbar(self):
        self.toolbar = QToolBar(self)
        self.toolbar.setContentsMargins(0, 0, 0, 0)
        # self.toolbar.setStyleSheet("QToolBar { background-color: lightgray; border: 0px solid gray; }")

        # Navigation buttons
        self.first_action = QAction(QIcon("Icons/icons8-thick-arrow-pointing-up-50.png"), "First", self)
        self.first_action.triggered.connect(self.show_first_page)
        self.toolbar.addAction(self.first_action)

        self.prev_action = QAction(QIcon("Icons/icons8-left-arrow-50.png"), "Previous", self)
        self.prev_action.triggered.connect(self.show_previous_page)
        self.toolbar.addAction(self.prev_action)

        # Page entry
        self.page_entry = QLineEdit(self)
        self.page_entry.setMaximumWidth(60)
        self.page_entry.setAlignment(Qt.AlignRight)
        self.page_entry.returnPressed.connect(self.go_to_page)
        self.toolbar.addWidget(self.page_entry)

        # Total pages label
        self.total_pages_label = QLabel(f"of {len(self.pdf)}" if self.pdf else "of 0")
        self.toolbar.addWidget(self.total_pages_label)

        # Navigation buttons
        self.next_action = QAction(QIcon("Icons/icons8-right-arrow-50.png"), "Next", self)
        self.next_action.triggered.connect(self.show_next_page)
        self.toolbar.addAction(self.next_action)

        self.last_action = QAction(QIcon("Icons/icons8-down-50.png"), "Last", self)
        self.last_action.triggered.connect(self.show_last_page)
        self.toolbar.addAction(self.last_action)

        # add separator
        self.toolbar.addSeparator()

        # Zoom actions
        self.zoom_in_action = QAction(QIcon("Icons/icons8-zoom-in-50.png"), "Zoom In", self)
        self.zoom_in_action.triggered.connect(self.zoom_in)
        self.toolbar.addAction(self.zoom_in_action)

        # QLineEdit for zoom percentage
        self.zoom_percentage_entry = QLineEdit(self)
        self.zoom_percentage_entry.setFixedWidth(60)  # Set a fixed width for consistency
        self.zoom_percentage_entry.setAlignment(Qt.AlignRight)
        self.zoom_percentage_entry.setPlaceholderText("100%")  # Default zoom is 100%
        self.zoom_percentage_entry.returnPressed.connect(self.set_zoom_from_entry)
        self.toolbar.addWidget(self.zoom_percentage_entry)

        self.zoom_out_action = QAction(QIcon("Icons/icons8-zoom-out-50.png"), "Zoom Out", self)
        self.zoom_out_action.triggered.connect(self.zoom_out)
        self.toolbar.addAction(self.zoom_out_action)

        # Create a reset zoom button with its icon and add it to the toolbar
        reset_zoom_icon = QIcon("Icons/icons8-zoom-to-actual-size-50.png")  # Replace with your icon path
        self.reset_zoom_action = QAction(reset_zoom_icon, "Reset Zoom", self)
        self.reset_zoom_action.triggered.connect(self.reset_zoom)
        self.toolbar.addAction(self.reset_zoom_action)

        # add separator
        self.toolbar.addSeparator()

        # Fit in window
        fit_window_icon = QIcon("Icons/icons8-enlarge-50.png")  # Replace with your icon path
        self.fit_window_action = QAction(fit_window_icon, "Fit in Window", self)
        self.fit_window_action.triggered.connect(self.fit_window)
        self.toolbar.addAction(self.fit_window_action)

        # Fit in width
        fit_width_icon = QIcon("Icons/icons8-resize-horizontal-50.png")  # Replace with your icon path
        self.fit_width_action = QAction(fit_width_icon, "Fit in Width", self)
        self.fit_width_action.triggered.connect(self.fit_width)
        self.toolbar.addAction(self.fit_width_action)

        # add separator
        self.toolbar.addSeparator()

        # Rotate left
        rotate_left_icon = QIcon("Icons/icons8-rotate-left-50.png")  # Replace with your icon path
        self.rotate_left_action = QAction(rotate_left_icon, "Rotate Left", self)
        self.rotate_left_action.triggered.connect(self.rotate_left)
        self.toolbar.addAction(self.rotate_left_action)

        # Rotate right
        rotate_right_icon = QIcon("Icons/icons8-rotate-right-50.png")  # Replace with your icon path
        self.rotate_right_action = QAction(rotate_right_icon, "Rotate Right", self)
        self.rotate_right_action.triggered.connect(self.rotate_right)
        self.toolbar.addAction(self.rotate_right_action)

        # add separator
        self.toolbar.addSeparator()

        # Pan tool button
        self.pan_tool_icon = QIcon("Icons/icons8-drag-50.png")  # Replace with your pan icon path
        self.pan_tool_action = QAction(self.pan_tool_icon, "Pan Tool", self)
        self.pan_tool_action.setCheckable(True)
        self.pan_tool_action.toggled.connect(self.toggle_pan_mode)
        self.toolbar.addAction(self.pan_tool_action)

        # add separator
        self.toolbar.addSeparator()

        # Print button
        self.print_icon = QIcon("Icons/icons8-print-50.png")  # Replace with your print icon path
        self.print_action = QAction(self.print_icon, "Print", self)
        self.print_action.triggered.connect(self.print_pdf)
        self.toolbar.addAction(self.print_action)

        self.save_pdf_action = QAction(QIcon("Icons/icons8-save-as-50.png"), "Save PDF", self)
        self.save_pdf_action.triggered.connect(self.save_pdf)
        self.toolbar.addAction(self.save_pdf_action)

    def setup_pdf_display_area(self):
        self.page_label = QLabel(self)
        self.page_label.setContentsMargins(0, 0, 0, 0)
        self.page_label.setAlignment(Qt.AlignCenter)

        self.scroll_area = QScrollArea(self)
        self.scroll_area.setContentsMargins(0, 0, 0, 0)
        self.scroll_area.setWidget(self.page_label)
        self.scroll_area.setWidgetResizable(True)

    def set_current_page(self, page_num):
        if not self.pdf:
            return
        max_pages = len(self.pdf)
        if 0 <= page_num < max_pages:
            self.current_page = page_num
            self.show_page(page_num)

    def go_to_page(self):
        try:
            page_num = int(self.page_entry.text()) - 1  # Minus 1 because pages start from 0
            self.set_current_page(page_num)
        except ValueError:
            QMessageBox.warning(self, "Invalid Page Number", "Please enter a valid page number.")

    def update_navigation_states(self):
        if not self.pdf:
            self.prev_action.setEnabled(False)
            self.next_action.setEnabled(False)
            return

        self.prev_action.setEnabled(self.current_page > 0)
        self.next_action.setEnabled(self.current_page < len(self.pdf) - 1)
        self.total_pages_label.setText(f"of {len(self.pdf)}")
        self.page_entry.setText(str(self.current_page + 1))

    def show_previous_page(self):
        self.set_current_page(self.current_page - 1)
        self.update_navigation_states()

    def show_next_page(self):
        self.set_current_page(self.current_page + 1)
        self.update_navigation_states()

    def show_page(self, page_num):
        if not self.pdf:
            return
        try:
            page = self.pdf[page_num]
            mat = Matrix(self.zoom_factor, self.zoom_factor).prerotate(self.rotation_angle)
            image = page.get_pixmap(matrix=mat)

            qt_image = QImage(image.samples, image.width, image.height, image.stride, QImage.Format_RGB888)
            pixmap = QPixmap.fromImage(qt_image)
            self.page_label.setPixmap(pixmap)
            self.update_navigation_states()
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to render page: {e}")

    def display(self, content):
        if self.pdf:
            self.pdf.close()
            self.pdf = None

        if content:
            # Try multiple recovery methods
            self.pdf = self._try_load_pdf_with_recovery(content)
            
            if self.pdf:
                try:
                    self.current_page = 0
                    self.show_page(self.current_page)
                    self.update_navigation_states()
                except Exception as e:
                    print(f"Error displaying PDF page: {e}")
                    self.page_label.clear()
                    self.page_label.setText(f"PDF loaded but page display failed: {str(e)}")
            else:
                self.page_label.clear()
                self.page_label.setText("Failed to load PDF - file may be corrupted")
        else:
            self.page_label.clear()

    def _try_load_pdf_with_recovery(self, content):
        """Try multiple methods to load a PDF, including recovery options."""
        pdf = None
        
        # Method 1: Normal opening
        try:
            pdf = fitz_open(stream=content, filetype="pdf")
            if pdf and len(pdf) > 0:
                return pdf
            if pdf:
                pdf.close()
        except Exception as e:
            print(f"Normal PDF load failed: {e}")
        
        # Method 2: Try to reconstruct from temporary file (sometimes file-based opening works better)
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name
            
            try:
                # Try opening from file
                pdf = fitz_open(tmp_path)
                if pdf and len(pdf) > 0:
                    print("PDF loaded from file")
                    return pdf
                if pdf:
                    pdf.close()
            finally:
                # Clean up temp file
                try:
                    os.unlink(tmp_path)
                except:
                    pass
        except Exception as e:
            print(f"File-based PDF recovery failed: {e}")
        
        # Method 3: Try to extract and reconstruct pages
        try:
            pdf = self._reconstruct_corrupted_pdf(content)
            if pdf and len(pdf) > 0:
                print("PDF reconstructed from corrupted file")
                return pdf
            if pdf:
                pdf.close()
        except Exception as e:
            print(f"PDF reconstruction failed: {e}")
        
        # Method 4: Try alternative PDF library (PyPDF2) as last resort
        try:
            pdf = self._try_pypdf2_recovery(content)
            if pdf:
                return pdf
        except Exception as e:
            print(f"PyPDF2 recovery failed: {e}")
        
        # Method 5: Try to extract raw text/bytes and create a simple text PDF
        try:
            pdf = self._create_text_pdf_from_content(content)
            if pdf:
                return pdf
        except Exception as e:
            print(f"Text PDF creation failed: {e}")
        
        return None

    def _reconstruct_corrupted_pdf(self, content):
        """Attempt to reconstruct a PDF from corrupted content."""
        try:
            import fitz
            import io
            
            # Try to create a new PDF and copy pages from corrupted one
            new_pdf = fitz.open()  # Create empty PDF
            
            # Try to open corrupted PDF with maximum error tolerance
            try:
                # Try opening from stream
                corrupted_pdf = fitz_open(stream=content, filetype="pdf")
                
                # Try to extract each page
                for page_num in range(len(corrupted_pdf)):
                    try:
                        page = corrupted_pdf[page_num]
                        # Try to get pixmap (renders the page)
                        pix = page.get_pixmap(matrix=fitz.Matrix(1, 1))
                        if pix:
                            # Insert page as image into new PDF
                            img_rect = fitz.Rect(0, 0, pix.width, pix.height)
                            new_page = new_pdf.new_page(width=pix.width, height=pix.height)
                            new_page.insert_image(img_rect, pixmap=pix)
                    except Exception as e:
                        print(f"Failed to extract page {page_num}: {e}")
                        continue
                
                corrupted_pdf.close()
                
                if len(new_pdf) > 0:
                    return new_pdf
                else:
                    new_pdf.close()
            except:
                pass
            
            # If page extraction failed, try text extraction and create text-based PDF
            try:
                corrupted_pdf = fitz_open(stream=content, filetype="pdf")
                text_content = []
                
                for page_num in range(len(corrupted_pdf)):
                    try:
                        page = corrupted_pdf[page_num]
                        text = page.get_text()
                        if text:
                            text_content.append((page_num, text))
                    except:
                        continue
                
                corrupted_pdf.close()
                
                if text_content:
                    # Create a new PDF with extracted text
                    for page_num, text in text_content:
                        page = new_pdf.new_page()
                        # Insert text as text (not image)
                        try:
                            page.insert_text((50, 50), text[:1000], fontsize=12)  # Limit text length
                        except:
                            pass
                    
                    if len(new_pdf) > 0:
                        return new_pdf
            except:
                pass
            
            new_pdf.close()
        except Exception as e:
            print(f"PDF reconstruction error: {e}")
        
        return None

    def _try_pypdf2_recovery(self, content):
        """Try to recover PDF using PyPDF2 as alternative."""
        try:
            from PyPDF2 import PdfReader, PdfWriter
            import io
            
            # Try to read with PyPDF2
            pdf_stream = io.BytesIO(content)
            reader = PdfReader(pdf_stream, strict=False)  # strict=False allows some errors
            
            if len(reader.pages) > 0:
                # Create a new PDF with PyMuPDF from PyPDF2 pages
                import fitz
                new_pdf = fitz.open()
                
                for page_num in range(len(reader.pages)):
                    try:
                        page = reader.pages[page_num]
                        # Extract text
                        text = page.extract_text()
                        
                        # Create new page with text
                        new_page = new_pdf.new_page()
                        if text:
                            new_page.insert_text((50, 50), text[:2000], fontsize=11)
                    except:
                        continue
                
                if len(new_pdf) > 0:
                    return new_pdf
                new_pdf.close()
        except Exception as e:
            print(f"PyPDF2 recovery attempt failed: {e}")
        
        return None

    def _create_text_pdf_from_content(self, content):
        """Create a simple text-based PDF from corrupted PDF content."""
        try:
            import fitz
            import re
            
            # Try to extract readable text from PDF bytes
            text_content = []
            
            # Look for text streams in PDF (between BT and ET markers)
            content_str = content.decode('latin-1', errors='ignore')
            
            # Find text blocks
            text_pattern = r'BT\s*(.*?)\s*ET'
            matches = re.finditer(text_pattern, content_str, re.DOTALL)
            
            for match in matches:
                text_block = match.group(1)
                # Try to extract readable text
                readable_text = re.sub(r'[^\x20-\x7E\n]', '', text_block)  # Keep printable ASCII
                if readable_text.strip():
                    text_content.append(readable_text[:500])  # Limit length
            
            # If we found text, create a PDF
            if text_content:
                new_pdf = fitz.open()
                for text in text_content:
                    try:
                        page = new_pdf.new_page()
                        # Insert text in chunks to avoid overflow
                        lines = text.split('\n')[:30]  # Limit to 30 lines
                        y_pos = 50
                        for line in lines:
                            if line.strip():
                                page.insert_text((50, y_pos), line[:80], fontsize=10)
                                y_pos += 15
                    except:
                        continue
                
                if len(new_pdf) > 0:
                    print("Created text PDF from corrupted content")
                    return new_pdf
                new_pdf.close()
        except Exception as e:
            print(f"Text PDF creation error: {e}")
        
        return None

    def clear(self):
        if self.pdf:
            self.pdf.close()
            self.pdf = None
        self.page_label.clear()

    def show_first_page(self):
        if self.pdf:
            self.set_current_page(0)

    def show_last_page(self):
        if self.pdf:
            self.set_current_page(len(self.pdf) - 1)

    def zoom_in(self):
        if self.pdf:
            self.zoom_factor *= 1.2  # Assuming you have initialized zoom_factor as 1 in your __init__ method
            self.update_zoom()

    def zoom_out(self):
        if self.pdf:
            self.zoom_factor *= 0.8
            self.update_zoom()

    def update_zoom(self):
        if not self.pdf:
            return
        # Always zoom on the original high-quality image
        page = self.pdf[self.current_page]
        mat = Matrix(self.zoom_factor, self.zoom_factor)
        image = page.get_pixmap(matrix=mat)

        qt_image = QImage(image.samples, image.width, image.height, image.stride, QImage.Format_RGB888)
        pixmap = QPixmap.fromImage(qt_image)

        self.page_label.setPixmap(pixmap)

    def set_zoom_from_entry(self):
        try:
            # Extract the percentage from the QLineEdit
            percentage = float(self.zoom_percentage_entry.text().strip('%')) / 100
            print(f"Entered Percentage: {percentage}")  # Debug print statement

            if 0.1 <= percentage <= 5:  # Just to ensure reasonable zoom limits, you can adjust these values
                self.zoom_factor = percentage
                self.show_page(self.current_page)
            else:
                QMessageBox.warning(self, "Invalid Zoom", "Please enter a zoom percentage between 10% and 500%.")
        except ValueError:
            QMessageBox.warning(self, "Invalid Zoom", "Please enter a valid zoom percentage.")

    def reset_zoom(self):
        if self.pdf:
            self.zoom_factor = 1.0
            self.show_page(self.current_page)
            self.zoom_percentage_entry.setText("100")

    def fit_window(self):
        if not self.pdf:
            return
        page = self.pdf[self.current_page]
        zoom_x = self.scroll_area.width() / page.rect.width
        zoom_y = self.scroll_area.height() / page.rect.height
        self.zoom_factor = min(zoom_x, zoom_y)
        self.show_page(self.current_page)

    def fit_width(self):
        if not self.pdf:
            return
        page = self.pdf[self.current_page]
        self.zoom_factor = self.scroll_area.width() / page.rect.width
        self.show_page(self.current_page)

    def rotate_left(self):
        if self.pdf:
            self.rotation_angle -= 90
            self.show_page(self.current_page)

    def rotate_right(self):
        if self.pdf:
            self.rotation_angle += 90
            self.show_page(self.current_page)

    def mousePressEvent(self, event):
        if event.button() == Qt.LeftButton and self.pan_mode:
            # Set the is_panning flag to True
            self.is_panning = True
            # Store the initial mouse position
            self.pan_start_x = event.x()
            self.pan_start_y = event.y()
            # Change the cursor to an open hand symbol
            self.setCursor(Qt.OpenHandCursor)
        event.accept()

    def mouseMoveEvent(self, event):
        if self.is_panning and self.pan_mode:
            # Calculate the distance moved by the mouse
            dx = event.x() - self.pan_start_x
            dy = event.y() - self.pan_start_y

            # Scroll the QScrollArea accordingly
            self.scroll_area.horizontalScrollBar().setValue(self.scroll_area.horizontalScrollBar().value() - dx)
            self.scroll_area.verticalScrollBar().setValue(self.scroll_area.verticalScrollBar().value() - dy)

            # Update the initial mouse position for the next mouse move event
            self.pan_start_x = event.x()
            self.pan_start_y = event.y()
        event.accept()

    def mouseReleaseEvent(self, event):
        if event.button() == Qt.LeftButton and self.is_panning and self.pan_mode:
            # Reset the is_panning flag
            self.is_panning = False
            # Reset the cursor
            self.setCursor(Qt.ArrowCursor)
        event.accept()

    def toggle_pan_mode(self, checked):
        if checked:
            self.pan_mode = True
            self.setCursor(Qt.OpenHandCursor)
        else:
            self.pan_mode = False
            self.setCursor(Qt.ArrowCursor)

    def print_pdf(self):
        if not self.pdf:
            QMessageBox.warning(self, "No Document", "No document available to print.")
            return

        printer = QPrinter()
        printer.setFullPage(True)
        printer.setPageOrientation(QPageLayout.Portrait)

        print_dialog = QPrintDialog(printer, self)
        if print_dialog.exec_() == QPrintDialog.Accepted:
            from PySide6.QtGui import QPainter

            painter = QPainter()
            if not painter.begin(printer):
                return

            num_pages = len(self.pdf)
            for i in range(num_pages):
                if i != 0:  # start a new page after the first one
                    printer.newPage()
                page = self.pdf[i]
                image = page.get_pixmap()
                qt_image = QImage(image.samples, image.width, image.height, image.stride, QImage.Format_RGB888)
                pixmap = QPixmap.fromImage(qt_image)
                rect = painter.viewport()
                size = pixmap.size()
                size.scale(rect.size(), Qt.KeepAspectRatio)
                painter.setViewport(rect.x(), rect.y(), size.width(), size.height())
                painter.setWindow(pixmap.rect())
                painter.drawPixmap(0, 0, pixmap)

            painter.end()

    def save_pdf(self):
        if not self.pdf:
            QMessageBox.warning(self, "No Document", "No document available to save.")
            return

        options = QFileDialog.Options()
        filePath, _ = QFileDialog.getSaveFileName(self, "Save PDF", "", "PDF Files (*.pdf);;All Files (*)",
                                                  options=options)

        if not filePath:
            return  # user cancelled the dialog

        if not filePath.endswith(".pdf"):
            filePath += ".pdf"

        try:
            self.pdf.save(filePath)  # save the PDF to the specified path
            QMessageBox.information(self, "Success", "PDF saved successfully!")
        except Exception as e:
            QMessageBox.critical(self, "Error", f"Failed to save PDF: {e}")


class AudioVideoViewer(QWidget):
    def __init__(self, parent=None):
        super(AudioVideoViewer, self).__init__(parent)

        # # Initialize the volumes control interface once
        # devices = AudioUtilities.GetSpeakers()
        # self.volume_interface = devices.Activate(
        #     IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
        # self.volume = cast(self.volume_interface, POINTER(IAudioEndpointVolume))

        # Initialize the volumes control
        if os.name == 'nt':  # Windows
            devices = AudioUtilities.GetSpeakers()
            self.volume_interface = devices.Activate(
                IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
            self.volume = self.volume_interface.QueryInterface(IAudioEndpointVolume)

        self.layout = QVBoxLayout(self)

        self._audio_output = QAudioOutput()
        self._player = QMediaPlayer()

        self._player.setAudioOutput(self._audio_output)

        self._video_widget = QVideoWidget(self)
        self.layout.addWidget(self._video_widget)
        self._player.setVideoOutput(self._video_widget)

        # Progress layout
        self.progress_layout = QHBoxLayout()

        # Progress Slider
        self.progress_slider = QSlider(Qt.Horizontal, self)
        self.progress_slider.setToolTip("Progress")
        self.progress_slider.setRange(0, self._player.duration())
        self.progress_slider.sliderMoved.connect(self.set_media_position)
        self.progress_slider.mousePressEvent = self.slider_clicked
        self.progress_layout.addWidget(self.progress_slider)

        # Progress label
        self.progress_label = QLabel("00:00", self)
        self.progress_layout.addWidget(self.progress_label)

        self.layout.addLayout(self.progress_layout)

        # Controls layout
        self.controls_layout = QHBoxLayout()

        # Spacer to push media control buttons to the center
        self.controls_layout.addSpacerItem(QSpacerItem(20, 10, QSizePolicy.Expanding, QSizePolicy.Minimum))

        # If system is Windows, add a volume slider
        if os.name == 'nt':
            # Volume label
            self.controls_layout.addWidget(QLabel("Volume"))

            # Volume slider
            self.volume_slider = QSlider(Qt.Horizontal, self)
            self.volume_slider.setToolTip("Volume")
            self.volume_slider.setRange(0, 100)
            self.volume_slider.setValue(self.get_system_volume())
            self.volume_slider.setFixedWidth(150)
            self.volume_slider.valueChanged.connect(self.update_volume_display)
            self.volume_slider.valueChanged.connect(self.set_volume)
            self.controls_layout.addWidget(self.volume_slider)

            # Volume display label
            self.volume_display = QLabel(f"{self.get_system_volume()}%", self)
            self.volume_display.setToolTip("Volume Percentage")
            self.controls_layout.addWidget(self.volume_display)

        # Spacer to separate media controls and volumes controls
        self.controls_layout.addSpacerItem(QSpacerItem(370, 10, QSizePolicy.Fixed, QSizePolicy.Minimum))

        icon_size = QSize(24, 24)

        # Media control buttons
        self.play_btn = QPushButton(self)
        self.play_btn.setToolTip("Play")
        self.play_btn.setIcon(QIcon("Icons/icons8-circled-play-50.png"))
        self.play_btn.setIconSize(icon_size)
        self.play_btn.clicked.connect(self._player.play)
        self.controls_layout.addWidget(self.play_btn)

        self.pause_btn = QPushButton(self)
        self.pause_btn.setToolTip("Pause")
        self.pause_btn.setIcon(QIcon("Icons/icons8-pause-button-50.png"))
        self.pause_btn.setIconSize(icon_size)
        self.pause_btn.clicked.connect(self._player.pause)
        self.controls_layout.addWidget(self.pause_btn)

        self.stop_btn = QPushButton(self)
        self.stop_btn.setToolTip("Stop")
        self.stop_btn.setIcon(QIcon("Icons/icons8-stop-circled-50.png"))
        self.stop_btn.setIconSize(icon_size)
        self.stop_btn.clicked.connect(self._player.stop)
        self.controls_layout.addWidget(self.stop_btn)

        # Spacer to separate volumes controls and speed controls
        self.controls_layout.addSpacerItem(QSpacerItem(370, 10, QSizePolicy.Fixed, QSizePolicy.Minimum))

        # Speed label
        self.controls_layout.addWidget(QLabel("Speed"))

        # Playback speed dropdown
        self.playback_speed_combo = QComboBox(self)
        self.playback_speed_combo.setToolTip("Playback Speed")
        speeds = ["0.25x", "0.5x", "0.75x", "1.0x", "1.25x", "1.5x", "1.75x", "2.0x"]
        self.playback_speed_combo.addItems(speeds)
        self.playback_speed_combo.setCurrentText("1.0x")
        self.playback_speed_combo.currentTextChanged.connect(self.change_playback_speed)
        self.controls_layout.addWidget(self.playback_speed_combo)

        # Spacer to push speed controls to the right
        self.controls_layout.addSpacerItem(QSpacerItem(20, 10, QSizePolicy.Expanding, QSizePolicy.Minimum))

        self.layout.addLayout(self.controls_layout)

        self._player.positionChanged.connect(self.update_position)
        self._player.positionChanged.connect(self.update_slider_position)
        self._player.durationChanged.connect(self.update_duration)

    def display(self, content, file_type="video"):
        self.playback_speed_combo.setCurrentText("1.0x")
        self._player.setPlaybackRate(1.0)
        self._player.setSource(QUrl.fromLocalFile(content))
        # very_old# self._player.play()

    def update_position(self, position):
        self.progress_label.setText("{:02d}:{:02d}".format(position // 60000, (position // 1000) % 60))

    def update_duration(self, duration):
        self.progress_slider.setRange(0, duration)
        self.progress_label.setText("{:02d}:{:02d} / {:02d}:{:02d}".format(self._player.position() // 60000,
                                                                           (self._player.position() // 1000) % 60,
                                                                           duration // 60000,
                                                                           (duration // 1000) % 60))

    def clear(self):
        self._player.stop()
        # clear the thumbnail of the video or audio file

    def change_playback_speed(self, speed_text):
        speed = float(speed_text.replace("x", ""))
        self._player.setPlaybackRate(speed)

    def update_slider_position(self, position):
        self.progress_slider.setValue(position)

    def set_media_position(self, position):
        self._player.setPosition(position)

    def slider_clicked(self, event):
        # Update the slider position when clicked
        new_value = int(event.x() / self.progress_slider.width() * self.progress_slider.maximum())

        # Ensure the value is within range
        new_value = max(0, min(new_value, self.progress_slider.maximum()))

        self.progress_slider.setValue(new_value)
        self.set_media_position(new_value)

    # def get_system_volume(self):
    #     """Return the current system volumes as a value between 0 and 100."""
    #     current_volume = self.volume.GetMasterVolumeLevelScalar()
    #     return int(current_volume * 100)
    #
    # @Slot(int)
    # def set_volume(self, value):
    #     """Set the system volumes based on the slider's value."""
    #     self.volume.SetMasterVolumeLevelScalar(value / 100.0, None)

    # @Slot(int)
    # def update_volume_display(self, value):
    #     """Update the volumes display label based on the slider's value."""
    #     self.volume_display.setText(f"{value}%")

    def get_system_volume(self):
        """Return the current system volume as a value between 0 and 100."""
        if os.name == 'nt':  # Windows
            current_volume = self.volume.GetMasterVolumeLevelScalar()
            return int(current_volume * 100)

    @Slot(int)
    def set_volume(self, value):
        """Set the system volume based on the slider's value."""
        if os.name == 'nt':  # Windows
            self.volume.SetMasterVolumeLevelScalar(value / 100.0, None)

    @Slot(int)
    def update_volume_display(self, value):
        """Update the volumes display label based on the slider's value."""
        if os.name == 'nt':  # Only for Windows
            self.volume_display.setText(f"{value}%")

    @Slot(int)
    def set_position(self, position):
        """Set the position of the media playback based on the slider's position."""
        self._player.setPosition(position)

    @Slot(int)
    def update_slider_position(self, position):
        """Update the slider's position based on the media's playback position."""
        self.progress_slider.setValue(position)
